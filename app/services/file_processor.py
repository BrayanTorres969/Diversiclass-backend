from pathlib import Path
import io
from typing import Optional
from PyPDF2 import PdfReader
from docx import Document
from fastapi import UploadFile
from pptx import Presentation


async def extract_text_from_file(file: UploadFile, extension: str) -> str:
    """Extrae texto de archivos PDF, DOCX o PPTX."""
    try:
        content = await file.read()
        
        if extension == '.pdf':
            with io.BytesIO(content) as pdf_file:
                pdf_reader = PdfReader(pdf_file)
                text_pages = []
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():  # Verifica que haya texto
                        text_pages.append(page_text)
                if not text_pages:
                    raise ValueError("El PDF no contiene texto extraíble")
                return "\n".join(text_pages)
        
        elif extension == '.docx':
            with io.BytesIO(content) as doc_file:
                doc = Document(doc_file)
                paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
                if not paragraphs:
                    raise ValueError("El DOCX no contiene texto extraíble")
                return "\n".join(paragraphs)
        
        elif extension == '.pptx':
            with io.BytesIO(content) as ppt_file:
                prs = Presentation(ppt_file)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text and shape.text.strip():
                            text.append(shape.text)
                if not text:
                    raise ValueError("El PPTX no contiene texto extraíble")
                return "\n".join(text)
        
        raise ValueError(f"Formato de archivo no soportado: {extension}")
    
    except Exception as e:
        raise ValueError(f"Error al procesar archivo: {str(e)}") from e